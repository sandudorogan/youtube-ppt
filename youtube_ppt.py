#!/usr/bin/python

from pptx import Presentation
from pptx.util import Inches
from pytube import YouTube
from tqdm import tqdm
import argparse
import cv2
import numpy as np
import os
import shutil
import re


def delete_folder(folder_path):
    if os.path.exists(folder_path):
        shutil.rmtree(folder_path)


def parse_arguments():
    parser = argparse.ArgumentParser(description="Create a PowerPoint presentation from a YouTube video.")
    parser.add_argument("url", help="YouTube video URL")
    parser.add_argument("--crop", help="Crop rectangle in format x,y,width,height")
    parser.add_argument("--start", help="Start time in format MM:SS", default="00:00")
    parser.add_argument("--end", help="End time in format MM:SS", default=None)
    parser.add_argument("--output", help="Output PowerPoint file path", default=None)
    parser.add_argument('--no-cache', action='store_true', help="Disable caching and force redownload of video and regeneration of images.")
    return parser.parse_args()


def extract_video_id(url):
    match = re.search(r"v=([^&]+)", url)
    if match:
        return match.group(1)
    else:
        raise ValueError(f"Could not extract video ID from URL: {url}")


def download_video(url):
    yt = YouTube(url)
    video_id = f"{url.split('watch?v=')[-1]}.mp4"
    video_stream = yt.streams.filter(file_extension='mp4').get_highest_resolution()
    output_path = f"videos/{video_id}"

    if not os.path.exists("videos"):
        os.makedirs("videos")

    video_stream.download(output_path="videos", filename=video_id)

    return output_path


def mse(image1, image2):
    err = np.sum((image1.astype('float') - image2.astype('float')) ** 2)
    err /= float(image1.shape[0] * image1.shape[1])

    return err


def crop_speaker(frame, crop_rect=None):
    if crop_rect is None:
        return frame

    x, y, w, h = crop_rect
    cropped_frame = frame[y:y+h, x:x+w]

    return cropped_frame


def save_images(frames, video_title, output_folder):
    os.makedirs(output_folder, exist_ok=True)

    with tqdm(total=len(frames), desc="Saving images", ncols=100) as pbar:
        for i, frame in enumerate(frames):
            cv2.imwrite(f'{output_folder}/frame_{i:03d}.png', frame)
            pbar.update(1)

    return output_folder


def process_video(video_filename, crop_rect=None, start_time='0:0', end_time=None):
    cap = cv2.VideoCapture(video_filename)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    fps = int(cap.get(cv2.CAP_PROP_FPS))

    start_minute, start_second = map(float, start_time.split(':'))
    if end_time is not None:
        end_minute, end_second = map(float, end_time.split(':'))
    else:
        end_minute, end_second = divmod(total_frames / fps, 60)

    start_frame = int((start_minute * 60 + start_second) * fps)
    end_frame = int((end_minute * 60 + end_second) * fps)

    cap.set(cv2.CAP_PROP_POS_FRAMES, start_frame)

    ret, prev_frame = cap.read()
    if crop_rect:
        prev_frame = crop_speaker(prev_frame, crop_rect)

    unique_frames = [prev_frame]

    mse_threshold = 200

    with tqdm(total=end_frame - start_frame, desc="Processing video", ncols=100) as pbar:
        while cap.isOpened() and start_frame < end_frame:
            ret, current_frame = cap.read()
            pbar.update(1)

            if not ret:
                break

            if crop_rect:
                current_frame = crop_speaker(current_frame, crop_rect)

            if mse(current_frame, prev_frame) > mse_threshold:
                unique_frames.append(current_frame)
                prev_frame = current_frame

            start_frame += 1

    cap.release()

    return unique_frames


def create_pptx(image_folder, video_title, output_path=None):
    if output_path is None:
        output_path = image_folder

    pptx_filename = os.path.join(output_path, f"{video_title}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]

    image_files = [f for f in sorted(os.listdir(image_folder)) if f.endswith(".png")]

    for filename in tqdm(image_files, desc="Creating PowerPoint presentation", unit="image", ncols=100):
        slide = prs.slides.add_slide(blank_slide_layout)
        image_path = os.path.join(image_folder, filename)
        img = cv2.imread(image_path)

        height, width, _ = img.shape
        left, top, pic_width, pic_height = Inches(0), Inches(0), prs.slide_width, prs.slide_height

        slide.shapes.add_picture(image_path, left, top, width=pic_width, height=pic_height)

    prs.save(pptx_filename)
    return pptx_filename


def main():
    args = parse_arguments()
    crop_rect = tuple(map(int, args.crop.split(','))) if args.crop else None

    video_id = extract_video_id(args.url)
    video_filename = os.path.join("videos", f"{video_id}.mp4")
    crop_suffix = f"_crop_{'_'.join(map(str, crop_rect))}" if crop_rect else ""
    image_folder = os.path.join("images", f"{video_id}{crop_suffix}")
    pptx_filename = args.output if args.output else f"{video_id}.pptx"

    print(f"Processing video: {video_filename}")

    if args.no_cache:
        delete_folder(video_filename)
        delete_folder(image_folder)

    if not os.path.exists(video_filename):
        video_filename = download_video(args.url)

    if not os.path.exists(image_folder):
        os.makedirs(image_folder)
        unique_frames = process_video(video_filename, crop_rect, args.start, args.end)
        save_images(unique_frames, video_id, image_folder)
    else:
        print("Using cached images...")

    pptx_filename = create_pptx(image_folder, video_id, args.output)

    print(f"PowerPoint presentation created: {pptx_filename}")


if __name__ == "__main__":
    main()
